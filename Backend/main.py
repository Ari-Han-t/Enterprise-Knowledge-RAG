"""
Enterprise Knowledge RAG 2025
- FastAPI + ChromaDB
- Multimodal ingestion: PDF, PPTX, images, audio (Whisper)
- OCR fallback for scanned PDFs
- JWT auth, rate limiting
- Enterprise-only answers with citations and refusal on low relevance
"""

# ================= STANDARD LIBS =================
import os
import io
import re
import json
import uuid
import tempfile
from typing import List, Dict
from datetime import datetime, timedelta

# ================= FASTAPI =================
from fastapi import (
    FastAPI, File, UploadFile, HTTPException,
    Depends, Request
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials

# ================= LIMITERS =================
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.util import get_remote_address
from slowapi.errors import RateLimitExceeded

# ================= ENV / AUTH =================
from dotenv import load_dotenv
from jose import jwt
from passlib.context import CryptContext

# ================= LLM =================
from groq import Groq

# ================= FILE PROCESSING =================
from pypdf import PdfReader
from pdf2image import convert_from_bytes
from PIL import Image
from pptx import Presentation
import pytesseract

# ================= ML / VECTOR =================
from pydantic import BaseModel
from sentence_transformers import SentenceTransformer
import chromadb

# ================= CONFIG =================
load_dotenv()

APP_NAME = "Enterprise Knowledge RAG"

JWT_SECRET = os.getenv("JWT_SECRET", "change-me")
JWT_ALGO = "HS256"
JWT_EXP = timedelta(hours=24)

GROQ_API_KEY = os.getenv("GROQ_API_KEY")
GROQ_MODEL = os.getenv("GROQ_MODEL", "llama-3.3-70b-versatile")
if not GROQ_API_KEY:
    raise RuntimeError("GROQ_API_KEY not set")

client = Groq(api_key=GROQ_API_KEY)

pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_CMD", "/usr/bin/tesseract")
POPPLER_PATH = os.getenv("POPPLER_PATH", "/usr/bin")

CHROMA_DIR = os.getenv("CHROMA_DIR", "data/chroma")
USERS_FILE = "users.json"
BASE_DIR = "data/users"
EMBED_MODEL_NAME = os.getenv("EMBED_MODEL_NAME", "all-MiniLM-L6-v2")
MIN_SIMILARITY = float(os.getenv("MIN_SIMILARITY", "0.30"))
MAX_K = int(os.getenv("MAX_K", "6"))
WHISPER_MODEL = os.getenv("WHISPER_MODEL", "base")

embed_model = SentenceTransformer(EMBED_MODEL_NAME)
os.makedirs(CHROMA_DIR, exist_ok=True)
chroma_client = chromadb.PersistentClient(path=CHROMA_DIR)

_whisper_model = None


def get_whisper_model():
    global _whisper_model
    if _whisper_model is None:
        import whisper
        _whisper_model = whisper.load_model(WHISPER_MODEL)
    return _whisper_model


# ================= LIMITERS =================
def rate_limit_key(request: Request):
    auth = request.headers.get("authorization")
    if auth and auth.startswith("Bearer "):
        try:
            payload = jwt.decode(
                auth.split()[1],
                JWT_SECRET,
                algorithms=[JWT_ALGO]
            )
            return payload["sub"]
        except Exception:
            pass
    return get_remote_address(request)


limiter = Limiter(key_func=rate_limit_key)


# ================= AUTH =================
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
security = HTTPBearer()


def hash_password(pw: str):
    return pwd_context.hash(pw)


def verify_password(pw: str, hashed: str):
    return pwd_context.verify(pw, hashed)


def create_token(user_id: str):
    payload = {"sub": user_id, "exp": datetime.utcnow() + JWT_EXP}
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGO)


def get_current_user(
    creds: HTTPAuthorizationCredentials = Depends(security)
):
    try:
        payload = jwt.decode(creds.credentials, JWT_SECRET, algorithms=[JWT_ALGO])
        return payload["sub"]
    except Exception:
        raise HTTPException(401, "Invalid or expired token")


# ================= STORAGE =================
def user_dir(user_id: str):
    path = os.path.join(BASE_DIR, user_id)
    os.makedirs(path, exist_ok=True)
    return path


def paths(user_id: str):
    base = user_dir(user_id)
    return {
        "chat": f"{base}/chat.json"
    }


def load_users():
    return json.load(open(USERS_FILE)) if os.path.exists(USERS_FILE) else {}


def save_users(users):
    with open(USERS_FILE, "w", encoding="utf-8") as f:
        json.dump(users, f, indent=2)


def load_json(path):
    return json.load(open(path, encoding="utf-8")) if os.path.exists(path) else []


def save_json(path, data):
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)


def get_collection(user_id: str):
    name = f"enterprise_user_{user_id}"
    try:
        return chroma_client.get_collection(name)
    except Exception:
        return chroma_client.create_collection(name, metadata={"hnsw:space": "cosine"})


# ================= UTILS =================
def normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def split_sentences(text: str) -> List[str]:
    return [s for s in re.split(r"(?<=[.!?])\s+", text) if s]


def chunk_text(text: str, max_chars: int = 1200, overlap_chars: int = 200) -> List[str]:
    text = normalize_text(text)
    if not text:
        return []
    sentences = split_sentences(text)
    chunks = []
    current = ""
    for sentence in sentences:
        if len(current) + len(sentence) + 1 <= max_chars:
            current = f"{current} {sentence}".strip()
        else:
            if current:
                chunks.append(current)
            current = sentence
    if current:
        chunks.append(current)
    if overlap_chars > 0 and len(chunks) > 1:
        overlapped = []
        for i, chunk in enumerate(chunks):
            if i == 0:
                overlapped.append(chunk)
                continue
            prev = chunks[i - 1]
            overlap = prev[-overlap_chars:]
            overlapped.append(f"{overlap} {chunk}".strip())
        return overlapped
    return chunks


def embed_texts(texts: List[str]) -> List[List[float]]:
    vectors = embed_model.encode(texts, normalize_embeddings=True)
    return [v.tolist() for v in vectors]


# ================= EXTRACTION =================
def extract_pdf(b: bytes) -> str:
    reader = PdfReader(io.BytesIO(b))
    text = "\n".join(p.extract_text() or "" for p in reader.pages)
    if text.strip():
        return text
    images = convert_from_bytes(b, poppler_path=POPPLER_PATH)
    return "\n".join(pytesseract.image_to_string(img) for img in images)


def extract_pptx(b: bytes) -> str:
    prs = Presentation(io.BytesIO(b))
    return "\n".join(
        shape.text for slide in prs.slides
        for shape in slide.shapes if hasattr(shape, "text")
    )


def extract_image(b: bytes) -> str:
    return pytesseract.image_to_string(Image.open(io.BytesIO(b)))


def extract_audio(b: bytes) -> str:
    model = get_whisper_model()
    with tempfile.NamedTemporaryFile(delete=False, suffix=".audio") as tmp:
        tmp.write(b)
        tmp_path = tmp.name
    try:
        result = model.transcribe(tmp_path)
        return result.get("text", "")
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass


def extract_text_for_file(filename: str, data: bytes) -> str:
    name = filename.lower()
    if name.endswith(".txt"):
        return data.decode("utf-8", errors="ignore")
    if name.endswith(".pdf"):
        return extract_pdf(data)
    if name.endswith(".pptx"):
        return extract_pptx(data)
    if name.endswith((".png", ".jpg", ".jpeg")):
        return extract_image(data)
    if name.endswith((".mp3", ".wav", ".m4a", ".flac", ".ogg")):
        return extract_audio(data)
    return ""


def build_context(hits: List[Dict]) -> str:
    parts = []
    for h in hits:
        meta = h["metadata"]
        source = meta.get("source", "unknown")
        chunk_id = meta.get("chunk_id", "?")
        parts.append(f"Source: {source} | Chunk: {chunk_id}\n{h['text']}")
    return "\n\n".join(parts)


def refusal_message() -> str:
    return (
        "I do not have enough enterprise data to answer that. "
        "Please provide relevant documents."
    )


# ================= APP =================
app = FastAPI(title=APP_NAME)
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


# ================= MODELS =================
class AuthReq(BaseModel):
    email: str
    password: str


class AskReq(BaseModel):
    question: str
    k: int = 4


# ================= HEALTH =================
@app.get("/health")
def health():
    return {"status": "ok"}


# ================= AUTH ROUTES =================
@app.post("/auth/signup")
@limiter.limit("5/minute")
def signup(request: Request, req: AuthReq):
    users = load_users()
    if req.email in users:
        raise HTTPException(400, "User exists")

    user_id = str(uuid.uuid4())
    users[req.email] = {
        "id": user_id,
        "password": hash_password(req.password)
    }
    save_users(users)

    return {"access_token": create_token(user_id)}


@app.post("/auth/login")
@limiter.limit("10/minute")
def login(request: Request, req: AuthReq):
    users = load_users()
    user = users.get(req.email)

    if not user or not verify_password(req.password, user["password"]):
        raise HTTPException(401, "Invalid credentials")

    return {"access_token": create_token(user["id"])}


# ================= UPLOAD =================
@app.post("/upload")
@limiter.limit("10/hour")
async def upload(
    request: Request,
    files: List[UploadFile] = File(...),
    user_id: str = Depends(get_current_user)
):
    collection = get_collection(user_id)

    documents = []
    metadatas = []
    ids = []
    errors = []

    for file in files:
        data = await file.read()
        text = extract_text_for_file(file.filename, data)
        if not text.strip():
            errors.append({"file": file.filename, "error": "No extractable text"})
            continue

        chunks = chunk_text(text)
        if not chunks:
            errors.append({"file": file.filename, "error": "No chunks generated"})
            continue

        embeddings = embed_texts(chunks)
        file_id = str(uuid.uuid4())
        for i, chunk in enumerate(chunks):
            ids.append(f"{file_id}_{i}")
            documents.append(chunk)
            metadatas.append({
                "source": file.filename,
                "chunk_id": i,
                "file_id": file_id
            })

        collection.add(
            ids=ids[-len(chunks):],
            documents=documents[-len(chunks):],
            metadatas=metadatas[-len(chunks):],
            embeddings=embeddings
        )

    if not documents and errors:
        raise HTTPException(400, {"error": "No content extracted", "details": errors})

    return {
        "chunks_added": len(documents),
        "files_processed": len(files),
        "errors": errors
    }


# ================= ASK =================
@app.post("/ask")
@limiter.limit("30/minute")
def ask(
    request: Request,
    req: AskReq,
    user_id: str = Depends(get_current_user)
):
    collection = get_collection(user_id)
    if collection.count() == 0:
        return {
            "answer": refusal_message(),
            "sources": [],
            "matches": []
        }
    k = min(req.k, MAX_K)
    query_vec = embed_texts([req.question])[0]

    results = collection.query(
        query_embeddings=[query_vec],
        n_results=k,
        include=["documents", "metadatas", "distances"]
    )

    docs = results.get("documents", [[]])[0]
    metas = results.get("metadatas", [[]])[0]
    dists = results.get("distances", [[]])[0]

    hits = []
    for doc, meta, dist in zip(docs, metas, dists):
        if doc is None:
            continue
        similarity = 1.0 - float(dist) if dist is not None else 0.0
        if similarity >= MIN_SIMILARITY:
            hits.append({
                "text": doc,
                "metadata": meta or {},
                "similarity": round(similarity, 4)
            })

    if not hits:
        return {
            "answer": refusal_message(),
            "sources": [],
            "matches": []
        }

    context = build_context(hits)
    prompt = (
        "You are an enterprise knowledge assistant. "
        "Use ONLY the provided context. "
        "If the answer is not in the context, say you do not have enough enterprise data. "
        "Provide citations in square brackets with the source filename.\n\n"
        f"Context:\n{context}\n\nQuestion:\n{req.question}"
    )

    res = client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[
            {"role": "system", "content": "Answer strictly from context."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.1
    )

    answer = res.choices[0].message.content

    chat = load_json(paths(user_id)["chat"])
    chat.append({"q": req.question, "a": answer, "ts": datetime.utcnow().isoformat()})
    save_json(paths(user_id)["chat"], chat)

    sources = sorted({h["metadata"].get("source", "unknown") for h in hits})

    return {
        "answer": answer,
        "sources": sources,
        "matches": hits
    }


# ================= STREAM =================
@app.post("/ask/stream")
@limiter.limit("20/minute")
def ask_stream(
    request: Request,
    req: AskReq,
    user_id: str = Depends(get_current_user)
):
    collection = get_collection(user_id)
    if collection.count() == 0:
        def refusal_gen():
            yield f"data: {refusal_message()}\n\n"
            yield "data: [DONE]\n\n"

        return StreamingResponse(refusal_gen(), media_type="text/event-stream")
    k = min(req.k, MAX_K)
    query_vec = embed_texts([req.question])[0]

    results = collection.query(
        query_embeddings=[query_vec],
        n_results=k,
        include=["documents", "metadatas", "distances"]
    )

    docs = results.get("documents", [[]])[0]
    metas = results.get("metadatas", [[]])[0]
    dists = results.get("distances", [[]])[0]

    hits = []
    for doc, meta, dist in zip(docs, metas, dists):
        similarity = 1.0 - float(dist) if dist is not None else 0.0
        if doc and similarity >= MIN_SIMILARITY:
            hits.append({
                "text": doc,
                "metadata": meta or {},
                "similarity": round(similarity, 4)
            })

    if not hits:
        def refusal_gen():
            yield f"data: {refusal_message()}\n\n"
            yield "data: [DONE]\n\n"

        return StreamingResponse(refusal_gen(), media_type="text/event-stream")

    context = build_context(hits)
    prompt = (
        "You are an enterprise knowledge assistant. "
        "Use ONLY the provided context. "
        "If the answer is not in the context, say you do not have enough enterprise data. "
        "Provide citations in square brackets with the source filename.\n\n"
        f"Context:\n{context}\n\nQuestion:\n{req.question}"
    )

    def gen():
        full = ""
        stream = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[
                {"role": "system", "content": "Answer strictly from context."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.1,
            stream=True
        )

        for c in stream:
            if c.choices and c.choices[0].delta.content:
                t = c.choices[0].delta.content
                full += t
                yield f"data: {t}\n\n"

        chat = load_json(paths(user_id)["chat"])
        chat.append({"q": req.question, "a": full, "ts": datetime.utcnow().isoformat()})
        save_json(paths(user_id)["chat"], chat)

        yield "data: [DONE]\n\n"

    return StreamingResponse(gen(), media_type="text/event-stream")


# ================= HISTORY =================
@app.get("/history")
def history(user_id: str = Depends(get_current_user)):
    return load_json(paths(user_id)["chat"])


# ================= ERROR =================
@app.exception_handler(Exception)
async def error_handler(req: Request, exc: Exception):
    return JSONResponse(
        status_code=500,
        content={"error": str(exc)}
    )
